from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def inspect_pptx(path):
    try:
        prs = Presentation(path)
        print(f"File: {path}")
        print(f"Slide dimensions: {prs.slide_width.inches:.2f}x{prs.slide_height.inches:.2f} inches")
        
        for i, slide in enumerate(prs.slides):
            print(f"\n--- Slide {i+1} ---")
            for shape in slide.shapes:
                print(f"Shape: '{shape.name}', Type: {shape.shape_type}")
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs)
                        if not text.strip(): continue
                        print(f"  Text: '{text}'")
                        for run in paragraph.runs:
                            if not run.text.strip(): continue
                            font = run.font
                            print(f"    Run: '{run.text}'")
                            print(f"      Font: {font.name}, Size: {font.size.pt if font.size else 'N/A'}, Bold: {font.bold}")
                
                # Check for rectangle-like shapes (often AUTO_SHAPE type 1)
                if shape.shape_type == 1: # MSO_SHAPE_TYPE.AUTO_SHAPE
                    print(f"  AutoShape: x={shape.left.inches:.2f}, y={shape.top.inches:.2f}, w={shape.width.inches:.2f}, h={shape.height.inches:.2f}")

    except Exception as e:
        print(f"Error: {e}")

inspect_pptx("../../../../../Downloads/Layer Chart Examples(1).pptx")
