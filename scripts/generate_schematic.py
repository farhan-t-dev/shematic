import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import re

# ─── Color Palette ────────────────────────────────────────────────────────────
PALETTE = [
    "1B3A6B",  # deep navy        — A
    "B03A2E",  # crimson           — B
    "2E86C1",  # bright blue       — C
    "1E8449",  # forest green      — D
    "884EA0",  # medium purple     — E
    "E67E22",  # warm orange       — Kinsale
    "117A65",  # teal              — F
    "7B241C",  # dark red          — G
    "1A5276",  # dark steel        — H
    "6C3483",  # deep violet       — I
    "1D6A39",  # dark green        — J
    "626567",  # charcoal          — K
    "943126",  # rust red          — L
    "2E4057",  # midnight          — M
    "0E6655",  # jade              — N
    "A04000",  # burnt sienna      — O
]

def hex_to_rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def set_vert_anchor(txBox, anchor="ctr"):
    sp = txBox._element
    txBody = sp.find(qn("p:txBody"))
    if txBody is not None:
        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", anchor)

# ─── Parse Excel ──────────────────────────────────────────────────────────────
def find_layer_label(row):
    """Return layer label string from any column in the row, or None."""
    for cell in row:
        if cell is None:
            continue
        s = str(cell).strip()
        if s.startswith("$") and ("xs" in s or "Primary" in s):
            return s
    return None

def detect_columns(rows):
    """
    Auto-detect which columns hold: carrier_name, status, authorization, pct, premium.
    Returns dict with keys: name_col, status_col, auth_col, pct_col, prem_col
    Strategy: find the header row that contains 'Carrier' and 'Authorization' etc.
    """
    for row in rows:
        row_str = [str(c).lower().strip() if c else "" for c in row]
        if any("carrier" in c for c in row_str) and any("auth" in c for c in row_str):
            name_col   = next((i for i, c in enumerate(row_str) if "carrier" in c), 1)
            status_col = next((i for i, c in enumerate(row_str) if "status" in c), name_col + 1)
            auth_col   = next((i for i, c in enumerate(row_str) if "auth" in c), name_col + 2)
            pct_col    = next((i for i, c in enumerate(row_str) if "%" in c or "particip" in c), auth_col + 1)
            prem_col   = next((i for i, c in enumerate(row_str) if "share" in c and "prem" in c), pct_col + 1)
            if prem_col == pct_col:  # fallback
                prem_col = pct_col + 1
            return dict(name=name_col, status=status_col, auth=auth_col,
                        pct=pct_col, prem=prem_col)
    # Default fallback (original template layout)
    return dict(name=0, status=1, auth=2, pct=3, prem=4)

def parse_excel(path):
    wb = openpyxl.load_workbook(path, data_only=True)
    ws = wb.active
    rows = list(ws.values)

    # Account name: check col 0 and col 1 of first 3 rows
    account_name  = "Insurance Program Schematic"
    policy_period = ""
    HEADER_SKIP = {
        "share", "layer", "carrier", "status", "authorization", "total", "totals",
        "premium", "participation", "named insured", "as of", "rate", "fees",
        "comm", "comm %", "comm $", "ajg", "rps", "percent complete"
    }
    for row in rows[:4]:
        for ci in (0, 1):  # only check first two columns for account name
            cell = row[ci] if len(row) > ci else None
            if cell is None:
                continue
            s = str(cell).strip()
            if not s or s.lower() in HEADER_SKIP or len(s) < 4:
                continue
            if s.startswith("$"):
                continue
            if re.match(r'^\d{1,2}/\d{1,2}/\d{4}', s):
                if not policy_period:
                    policy_period = s
            elif account_name == "Insurance Program Schematic" and s != "Named Insured":
                account_name = s

    cols = detect_columns(rows)
    nc, sc, ac, pc, prc = cols["name"], cols["status"], cols["auth"], cols["pct"], cols["prem"]

    layers = []
    current_layer = None
    SKIP_NAMES = {"totals", "percent complete", "carrier", ""}

    for row in rows:
        # Pad row if needed
        row = list(row) + [None] * 10

        # Check any cell for a layer label
        lbl = find_layer_label(row)
        if lbl:
            current_layer = {"label": lbl, "carriers": []}
            layers.append(current_layer)
            continue

        if current_layer is None:
            continue

        name = str(row[nc]).strip() if row[nc] is not None else ""
        if name.lower() in SKIP_NAMES or name.startswith("$"):
            continue

        # Must have either auth or pct
        raw_auth = row[ac]
        raw_pct  = row[pc]
        if raw_auth is None and raw_pct is None:
            continue

        try:
            auth   = float(raw_auth) if raw_auth not in (None, "", "Included") else 0
            pct    = float(raw_pct)  if raw_pct  not in (None, "", "Included") else 0
            prem_v = row[prc]
            prem   = float(prem_v)   if prem_v   not in (None, "", "Included") else 0
            status = str(row[sc]).strip() if row[sc] else ""

            if auth > 0 or pct > 0:
                current_layer["carriers"].append({
                    "name": name, "auth": auth,
                    "pct": pct, "premium": prem, "status": status,
                })
        except (TypeError, ValueError):
            pass

    return account_name, policy_period, layers

def parse_layer_bounds(label):
    """Return (bottom, top) in dollars, or None if label is non-standard."""
    nums = [int(x.replace(",", "")) for x in re.findall(r'\$([0-9,]+)', label)]
    if not nums:
        return None
    if "Primary" in label:
        return 0, nums[0]
    if len(nums) >= 2:
        return nums[1], nums[1] + nums[0]
    return None

# ─── Build PPTX ──────────────────────────────────────────────────────────────
def build_schematic(excel_path, output_path, account_name_override=None, carrier_merges=None):
    account_name, policy_period, layers = parse_excel(excel_path)

    # Apply overrides
    if account_name_override:
        account_name = account_name_override
    
    if carrier_merges:
        for layer in layers:
            for c in layer["carriers"]:
                if c["name"] in carrier_merges:
                    c["name"] = carrier_merges[c["name"]]

    bounds_raw = [parse_layer_bounds(l["label"]) for l in layers]
    # Drop layers we can't parse bounds for (e.g. DBB non-standard)
    valid = [(l, b) for l, b in zip(layers, bounds_raw) if b is not None]
    # Clean each layer: drop zero-pct and zero-auth carriers
    for layer, _ in valid:
        real = [c for c in layer["carriers"] if c["pct"] > 0 and c["auth"] > 0]
        # Fall back to pct-only if no auth data at all (some sheets omit auth)
        if not real:
            real = [c for c in layer["carriers"] if c["pct"] > 0]
        layer["carriers"] = real
    # Drop layers with no carriers after cleaning
    valid = [(l, b) for l, b in valid if l["carriers"]]
    layers = [l for l, b in valid]
    bounds = [b for l, b in valid]
    program_max = max(b[1] for b in bounds)

    carrier_colors = {}
    ci = 0
    for layer in layers:
        for c in layer["carriers"]:
            if c["name"] not in carrier_colors:
                carrier_colors[c["name"]] = PALETTE[ci % len(PALETTE)]
                ci += 1

    prs = Presentation()
    prs.slide_width  = Inches(13.3)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb("F4F6F9")

    # Layout constants
    CHART_L = 1.40
    CHART_T = 1.00
    CHART_W = 11.50
    CHART_H = 6.20

    # ── Detect parallel layers (same bottom = interceding/side-by-side) ──
    from collections import defaultdict
    bottom_groups = defaultdict(list)
    for i, b in enumerate(bounds):
        bottom_groups[b[0]].append(i)

    # layer_x_frac[i] = (start_frac, width_frac) of CHART_W for that layer
    layer_x_frac = {}
    for bot, idxs in bottom_groups.items():
        n = len(idxs)
        for rank, i in enumerate(idxs):
            layer_x_frac[i] = (rank / n, 1.0 / n)

    # Y-axis: build using unique bottom values, using tallest layer at each bottom
    unique_bottoms = sorted(bottom_groups.keys())
    dominant_bounds = []  # one entry per unique bottom level
    for bot in unique_bottoms:
        idxs = bottom_groups[bot]
        tallest = max(idxs, key=lambda i: bounds[i][1])
        dominant_bounds.append(bounds[tallest])

    # Proportional height with minimum floor — based on dominant (tallest) bounds
    dom_sizes     = [b[1] - b[0] for b in dominant_bounds]
    total_dollars = sum(dom_sizes)
    MIN_FRAC      = 0.10
    raw_fracs     = [s / total_dollars for s in dom_sizes]
    boosted       = [max(f, MIN_FRAC) for f in raw_fracs]
    norm_fracs    = [f / sum(boosted) for f in boosted]

    sorted_dom = sorted(zip(dominant_bounds, norm_fracs), key=lambda x: x[0][0])

    y_map = {}
    y_cur = CHART_T + CHART_H
    for (bot, top), frac in sorted_dom:
        y_map[bot] = y_cur
        y_cur -= frac * CHART_H
        y_map[top] = y_cur
    y_map[program_max] = CHART_T

    # For parallel layers with different tops, interpolate their top within the band
    for i, b in enumerate(bounds):
        bot, top = b
        if top not in y_map:
            # Interpolate: find which dominant band this top falls in
            for (dbot, dtop), frac in sorted_dom:
                if dbot == bot:
                    band_h = y_map[dbot] - y_map[dtop]
                    ratio  = (top - dbot) / (dtop - dbot)
                    y_map[top] = y_map[dbot] - ratio * band_h
                    break

    def d2y(amount):
        if amount in y_map:
            return y_map[amount]
        for (bot, top), frac in sorted_dom:
            if bot < amount < top:
                lh = y_map[bot] - y_map[top]
                return y_map[bot] - ((amount - bot) / (top - bot)) * lh
        return CHART_T + CHART_H

    def add_rect(x, y, w, h, fill_hex, line_hex="FFFFFF", line_pt=0.75):
        s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        if line_hex:
            s.line.color.rgb = hex_to_rgb(line_hex)
            s.line.width = Pt(line_pt)
        else:
            s.line.fill.background()
        return s

    def add_textbox(x, y, w, h, lines, font_size=9,
                    align=PP_ALIGN.CENTER, valign="ctr",
                    rotate=0, wrap=True):
        """lines = list of (text, bold, color_hex, optional_size)"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.margin_left = tf.margin_right = Pt(3)
        tf.margin_top  = tf.margin_bottom = Pt(1)
        set_vert_anchor(tb, valign)

        if isinstance(lines, str):
            lines = [(lines, False, "333333", font_size)]

        first = True
        for item in lines:
            txt  = item[0]
            bld  = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else "333333"
            fsz  = item[3] if len(item) > 3 else font_size

            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            r = p.add_run()
            r.text = txt
            r.font.size  = Pt(fsz)
            r.font.bold  = bld
            r.font.color.rgb = hex_to_rgb(col)
            r.font.name  = "Calibri"

        if rotate:
            tb.rotation = rotate
        return tb

    # ── 1. Chart background ──
    add_rect(CHART_L, CHART_T, CHART_W, CHART_H, "FFFFFF", "BBBBBB", 0.5)

    # ── 2. Title ──
    add_textbox(0.1, 0.08, 13.1, 0.62,
                [(account_name, True, "1B3A6B", 24)],
                align=PP_ALIGN.CENTER, valign="ctr")

    if policy_period:
        add_textbox(0.1, 0.68, 13.1, 0.30,
                    [(policy_period, False, "666666", 11)],
                    align=PP_ALIGN.CENTER, valign="ctr")

    # ── 3. Y-axis dividers + labels (only at dominant boundary levels) ──
    # Use dominant bounds levels only to avoid duplicate lines from parallel layers
    dom_levels = sorted(set(v for b in dominant_bounds for v in b))

    for level in dom_levels:
        y = d2y(level)

        # Hairline rule across chart
        rule = slide.shapes.add_shape(
            1, Inches(CHART_L), Inches(y - 0.004),
            Inches(CHART_W), Inches(0.008))
        rule.fill.solid()
        rule.fill.fore_color.rgb = hex_to_rgb("999999")
        rule.line.fill.background()

        if level == 0:
            continue

        lbl = f"${level:,.0f}"
        is_bold = False
        add_textbox(0.0, y - 0.20, CHART_L - 0.05, 0.40,
                    [(lbl, is_bold, "222222", 10)],
                    align=PP_ALIGN.RIGHT, valign="ctr", wrap=False)

    # Also add lighter label for sub-tops of parallel layers (e.g. $15M when $20M is dominant)
    all_levels = sorted(set(v for b in bounds for v in b))
    for level in all_levels:
        if level not in dom_levels:
            y = d2y(level)
            rule = slide.shapes.add_shape(
                1, Inches(CHART_L), Inches(y - 0.003),
                Inches(CHART_W), Inches(0.006))
            rule.fill.solid()
            rule.fill.fore_color.rgb = hex_to_rgb("BBBBBB")
            rule.line.fill.background()
            
            if level == 0:
                continue

            lbl = f"${level:,.0f}"
            add_textbox(0.0, y - 0.18, CHART_L - 0.05, 0.36,
                        [(lbl, False, "222222", 9)],
                        align=PP_ALIGN.RIGHT, valign="ctr", wrap=False)

    # ── 4. Carrier blocks ──
    GAP = 0.016

    for li, (layer, (bot, top)) in enumerate(zip(layers, bounds)):
        y_top = d2y(top)
        y_bot = d2y(bot)
        blk_h = y_bot - y_top

        # Parallel layer x-offset
        x_start_frac, x_width_frac = layer_x_frac[li]
        layer_x_start = CHART_L + x_start_frac * CHART_W
        layer_x_width = x_width_frac * CHART_W

        # Vertical divider between parallel layers (not for first)
        if x_start_frac > 0:
            div = slide.shapes.add_shape(
                1, Inches(layer_x_start - 0.015), Inches(y_top),
                Inches(0.015), Inches(blk_h))
            div.fill.solid()
            div.fill.fore_color.rgb = hex_to_rgb("F4F6F9")
            div.line.fill.background()

        carriers  = layer["carriers"]
        tot_pct   = sum(c["pct"] for c in carriers) or 1.0
        x_cur     = layer_x_start

        for carrier in carriers:
            pct_n  = carrier["pct"] / tot_pct
            blk_w  = pct_n * layer_x_width
            color  = carrier_colors[carrier["name"]]

            dx = x_cur + GAP / 2
            dw = max(blk_w - GAP, 0.01)

            add_rect(dx, y_top, dw, blk_h, color, "FFFFFF", 0.6)

            # Label
            name       = carrier["name"]
            pct_val    = carrier["pct"] * 100
            pct_str    = f"{pct_val:.2f}%"
            dol_str    = f"${carrier['auth']:,.0f}"

            if dw < 0.22 or blk_h < 0.22:
                pass  # too small

            elif dw < 0.50:
                add_textbox(dx, y_top, dw, blk_h,
                            [(name, True, "FFFFFF", 7)],
                            align=PP_ALIGN.CENTER, valign="ctr",
                            rotate=270, wrap=False)

            elif dw < 0.85 or blk_h < 0.50:
                add_textbox(dx, y_top, dw, blk_h,
                            [(f"{name}\n{pct_str} ({dol_str})", True, "FFFFFF", 8)],
                            align=PP_ALIGN.CENTER, valign="ctr")

            else:
                add_textbox(dx, y_top, dw, blk_h, [
                    (name,    True,  "FFFFFF",  9),
                    (f"{pct_str} ({dol_str})", False, "E0E0E0",  8),
                ], align=PP_ALIGN.CENTER, valign="ctr")

            x_cur += blk_w

    prs.save(output_path)
    print(f"✅ Done → {output_path}")

if __name__ == "__main__":
    import sys
    if len(sys.argv) < 3:
        print("Usage: python generate_schematic.py input.xlsx output.pptx")
    else:
        build_schematic(sys.argv[1], sys.argv[2])
